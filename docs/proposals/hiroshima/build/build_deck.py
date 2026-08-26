# -*- coding: utf-8 -*-
"""広島市向け提案書ビルダー"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = '/root/.claude/uploads/33b70c79-5330-5515-8cdd-be08ae2c87f2/63841397-20260825____________________v11.pptx'
IMG  = '/tmp/w/img'
OUT  = '/tmp/w/out/20260827_広島市御中_機能拡張のご提案.pptx'

FONT   = 'Meiryo UI'
PRIMARY= RGBColor(0x00,0x81,0xAF)
DARK   = RGBColor(0x00,0x58,0x7A)
BORDER = RGBColor(0xD8,0xE2,0xE6)
BG1    = RGBColor(0xF7,0xFA,0xFB)
BG2    = RGBColor(0xF0,0xF5,0xF7)
TEXT   = RGBColor(0x26,0x26,0x26)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GRAY   = RGBColor(0x77,0x88,0x8E)
FLOWBD = RGBColor(0xA6,0xBC,0xC6)

L, CW = 0.30, 9.40          # 左マージン / コンテンツ幅

# ---------------------------------------------------------------- low level
def _rpr_font(run, size, color, bold=False):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', FONT)

def _bullet(par, on=True, size=10.5, color=PRIMARY):
    pPr = par._p.get_or_add_pPr()
    for t in ('a:buNone','a:buChar','a:buClr','a:buSzPct','a:buFont','a:lnSpc','a:spcAft','a:spcBef'):
        for e in pPr.findall(qn(t)): pPr.remove(e)
    def mk(tag, **kw):
        e = pPr.makeelement(qn(tag), kw); return e
    ln = mk('a:lnSpc'); p = ln.makeelement(qn('a:spcPct'), {'val':'112000'}); ln.append(p); pPr.append(ln)
    sb = mk('a:spcBef'); s = sb.makeelement(qn('a:spcPts'), {'val':'0'}); sb.append(s); pPr.append(sb)
    sa = mk('a:spcAft'); s = sa.makeelement(qn('a:spcPts'), {'val':'300'}); sa.append(s); pPr.append(sa)
    if on:
        ind = int(Pt(size).emu * 1.30); mar = ind + int(Pt(size).emu*0.22)
        pPr.set('marL', str(mar)); pPr.set('indent', str(-ind))
        bc = mk('a:buClr'); c = bc.makeelement(qn('a:srgbClr'), {'val':'%02X%02X%02X'%(color[0],color[1],color[2])}); bc.append(c); pPr.append(bc)
        pPr.append(mk('a:buSzPct', val='55000'))
        pPr.append(mk('a:buFont', typeface='Arial'))
        pPr.append(mk('a:buChar', char='●'))
    else:
        pPr.set('marL','0'); pPr.set('indent','0')
        pPr.append(mk('a:buNone'))

def textbox(slide, x, y, w, h, lines, size=11, color=TEXT, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, bullets=False, wrap=True):
    """lines: str または (str, dict) のリスト"""
    if isinstance(lines, str): lines = [lines]
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.margin_top = tf.margin_bottom = Emu(18288)
    for i, item in enumerate(lines):
        opt = {}
        if isinstance(item, tuple): item, opt = item
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = opt.get('align', align)
        sz = opt.get('size', size)
        _bullet(par, on=opt.get('bullets', bullets), size=sz, color=opt.get('bucolor', PRIMARY))
        r = par.add_run(); r.text = item
        _rpr_font(r, sz, opt.get('color', color), opt.get('bold', bold))
    return tb

def box(slide, x, y, w, h, lines, size=10.5, color=TEXT, fill=None, line=BORDER,
        bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, bullets=False,
        shape=MSO_SHAPE.RECTANGLE):
    if isinstance(lines, str): lines = [lines]
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(0.5)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(64008)
    tf.margin_top = tf.margin_bottom = Emu(36576)
    for i, item in enumerate(lines):
        opt = {}
        if isinstance(item, tuple): item, opt = item
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = opt.get('align', align)
        sz = opt.get('size', size)
        _bullet(par, on=opt.get('bullets', bullets), size=sz, color=opt.get('bucolor', PRIMARY))
        r = par.add_run(); r.text = item
        _rpr_font(r, sz, opt.get('color', color), opt.get('bold', bold))
    return sp

def arrow(slide, x, y, w=0.16, h=0.22, color=PRIMARY):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    return sp

# ---------------------------------------------------------------- slide parts
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[0])

def title(slide, text):
    textbox(slide, L, 0.15, 9.70, 0.38, text, size=20, color=RGBColor(0,0,0),
            anchor=MSO_ANCHOR.MIDDLE)

def lead(slide, text):
    textbox(slide, L, 0.63, CW, 0.42, text, size=12, color=TEXT)

def note(slide, text, y=6.92):
    textbox(slide, L, y, CW, 0.30, text, size=9, color=GRAY)

def secthead(slide, x, y, w, text, size=11):
    return box(slide, x, y, w, 0.34, text, size=size, color=WHITE, fill=PRIMARY,
               line=BORDER, align=PP_ALIGN.CENTER)

def label(slide, x, y, w, h, text, size=11):
    return box(slide, x, y, w, h, text, size=size, color=WHITE, fill=DARK,
               line=BORDER, align=PP_ALIGN.CENTER)

def grid(slide, x, y, widths, header, rows, hh=0.34, heights=None, size=10.5,
         first_col_label=True):
    """header: list[str] / rows: list[list[str]] （セル内は \n で箇条書き）"""
    cx = x
    for w, t in zip(widths, header):
        secthead(slide, cx, y, w, t); cx += w
    cy = y + hh
    for ri, row in enumerate(rows):
        rh = (heights[ri] if heights else 0.60)
        cx = x
        for ci, (w, cell) in enumerate(zip(widths, row)):
            first = (ci == 0 and first_col_label)
            lines = [l for l in str(cell).split('\n') if l != '']
            multi = len(lines) > 1
            box(slide, cx, cy, w, rh, lines if lines else '',
                size=size,
                color=DARK if first else TEXT,
                fill=BG2 if first else (BG1 if ri % 2 else WHITE),
                bold=first,
                align=PP_ALIGN.CENTER if first else PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE if (first or not multi) else MSO_ANCHOR.TOP,
                bullets=(not first) and multi)
            cx += w
        cy += rh
    return cy

def picture(slide, path, x, y, h=None, w=None, caption=None, cap_w=None):
    kw = {}
    if h: kw['height'] = Inches(h)
    if w: kw['width'] = Inches(w)
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)
    pic.line.color.rgb = BORDER; pic.line.width = Pt(0.75)
    if caption:
        cw = cap_w or pic.width / 914400
        textbox(slide, x, y + pic.height / 914400 + 0.04, cw, 0.24, caption,
                size=9, color=DARK, align=PP_ALIGN.CENTER)
    return pic

def pagenum(slide, n):
    textbox(slide, 9.05, 7.05, 0.70, 0.26, str(n), size=9, color=GRAY,
            align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------- 既存スライド流用
def para_text(par):
    return ''.join(r.text for r in par.runs)

def replace_paragraph(slide, mapping):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for par in sh.text_frame.paragraphs:
            t = para_text(par)
            if t in mapping:
                new = mapping[t]
                runs = par.runs
                if not runs: continue
                runs[0].text = new
                for r in runs[1:]:
                    r._r.getparent().remove(r._r)

def clone_shapes(src_slide, dst_slide):
    for sh in src_slide.shapes:
        dst_slide.shapes._spTree.append(copy.deepcopy(sh._element))

# ================================================================ 構築
prs = Presentation(BASE)
src = list(prs.slides)
SRC_SCOPE, SRC_FLOW_ABD, SRC_FLOW_C = src[3], src[5], src[6]
keep = {id(SRC_SCOPE): None, id(SRC_FLOW_ABD): None, id(SRC_FLOW_C): None}
# 既存スライドの中身を退避
stash = {}
for k, s in (('scope', SRC_SCOPE), ('abd', SRC_FLOW_ABD), ('c', SRC_FLOW_C)):
    stash[k] = [copy.deepcopy(sh._element) for sh in s.shapes]

# 既存スライドを全削除
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

def paste(slide, key):
    for el in stash[key]:
        slide.shapes._spTree.append(copy.deepcopy(el))

# ---------------------------------------------------------------- P1 表紙
s = new_slide(prs)
textbox(s, L, 1.20, 6.00, 0.45, '広島市 御中', size=16, color=DARK)
textbox(s, 0.28, 2.10, 9.45, 1.70,
        ['広島市モード・電子申請・広報・防災',
         '機能拡張のご提案'], size=32, color=RGBColor(0x40,0x40,0x40),
        anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.30, 3.95, 2.20, 0.05, '', fill=PRIMARY, line=None)
textbox(s, 0.28, 4.15, 8.50, 0.50,
        '〜 としポ基盤を活用した、広島市統合ポータルの実現 〜', size=16, color=DARK)
textbox(s, 1.83, 5.32, 7.72, 0.80,
        ['2026年8月27日', 'フェリカポケットマーケティング株式会社'],
        size=20, color=RGBColor(0x40,0x40,0x40), align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- P2 目次
s = new_slide(prs); title(s, '目次')
lead(s, '本資料は、営業ヒアリングを踏まえた4要件のご提案内容と、防災機能の中長期ロードマップで構成しています。')
toc = [('1．ご提案の背景と目的','P3'),
       ('2．ご提案の全体像','P4'),
       ('3．4要件の全体整理','P5'),
       ('4．機能概要と利用フロー','P6〜P8'),
       ('5．画面イメージ（アプリ／管理画面）','P9〜P16'),
       ('6．防災機能の拡張ロードマップ','P17〜P20'),
       ('7．運用・分析・システム基盤','P21〜P23'),
       ('8．実施スケジュールと今後の進め方','P24〜P25'),
       ('　　別紙一覧','P26')]
y = 1.30
for t, p in toc:
    box(s, L, y, 7.20, 0.44, t, size=12, color=TEXT, fill=BG1, line=BORDER,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 7.50, y, 2.20, 0.44, p, size=11, color=DARK, fill=BG2, line=BORDER,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.50
textbox(s, L, y + 0.10, CW, 0.60,
        ['別紙：アプリワイヤーフレーム_draft1_20260824.pdf ／ 管理画面ワイヤーフレーム_draft1_20260824.pdf ／ 20260824_現状機能概要まとめ.xlsx',
         '※「休日保育予約」は別資料にて整理のため、本資料の対象外としています。'],
        size=9.5, color=GRAY)

# ---------------------------------------------------------------- P3 背景と目的
s = new_slide(prs); title(s, '1．ご提案の背景と目的')
lead(s, '広島市が保有する既存資産を最大限に活かし、短期間・低リスクで統合ポータルの第一歩を実現します。')
label(s, L, 1.25, 1.70, 1.05, '背景')
textbox(s, 2.10, 1.25, 7.60, 1.05,
        ['市民向けサービス・情報が複数のサイト・アプリに分散し、必要な情報に到達するまでの導線が長い',
         '広島広域都市圏ポイント「としポ」は弊社開発アプリであり、事業者間調整なく最小限の負担で拡張が可能'],
        size=11, bullets=True)
label(s, L, 2.45, 1.70, 1.40, 'ご提案の目的')
textbox(s, 2.10, 2.45, 7.60, 1.40,
        ['「広島市モード」を新設し、市のサービス・情報をとしポアプリ上に一元集約する',
         '電子申請・広報・防災は既存資産への導線提供を基本とし、初年度は確実にリリースできる範囲に絞る',
         'デジタル庁のデジタル地方創生サービスカタログ掲載パッケージ「よむすび」を活用し、開発リスクを抑制する'],
        size=11, bullets=True)
label(s, L, 4.00, 1.70, 1.05, 'ご提供価値')
textbox(s, 2.10, 4.00, 7.60, 1.05,
        ['市民：手続・広報・防災の入口がひとつになり、必要な情報が必要なタイミングで届く',
         '広島市：アプリ改修を伴わない管理画面運用と、利用状況データに基づく施策改善が可能となる'],
        size=11, bullets=True)
textbox(s, L, 5.25, 4.00, 0.36, '今後の進め方', size=12, color=DARK, bold=True)
steps = ['①提案内容の\nご確認', '②要求仕様の\n確定', '③要件定義\n・設計', '④開発\n・テスト', '⑤リリース\n・運用']
x = 0.35
for i, st in enumerate(steps):
    box(s, x, 5.68, 1.62, 1.00, st.split('\n'), size=11, color=DARK, fill=BG2,
        line=BORDER, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1: arrow(s, x + 1.64, 6.07, 0.20, 0.22)
    x += 1.86
note(s, '※ 本資料の内容は、要件定義の結果により細部を調整させていただく場合があります。')

# ---------------------------------------------------------------- P4 全体像
s = new_slide(prs); title(s, '2．ご提案の全体像')
lead(s, 'としポアプリ内に「広島市モード」を新設し、市独自メニューと共通メニューを1画面に集約します。')
box(s, L, 1.22, CW, 0.34, '市民（アプリ利用者）', size=11, color=WHITE, fill=DARK,
    line=BORDER, align=PP_ALIGN.CENTER)
box(s, L, 1.62, CW, 0.62, '広島市モード ホーム（統合ポータル）　※居住地属性によりON／OFFを初期設定',
    size=12, color=WHITE, fill=PRIMARY, line=BORDER, align=PP_ALIGN.CENTER)
feat = [('A．広島市モード', 'マルチテナント方式で\n市独自メニューを表示'),
        ('B．電子申請', 'オンライン手続\nポータルサイトへ誘導'),
        ('C．広報', '広報紙PDFを\nアプリ内で閲覧'),
        ('D．防災', '「避難所へGo!」を起動\n未導入時はストアへ')]
x = L
for t, d in feat:
    box(s, x, 2.42, 2.26, 0.34, t, size=11, color=WHITE, fill=DARK, line=BORDER,
        align=PP_ALIGN.CENTER)
    box(s, x, 2.76, 2.26, 0.92, d.split('\n'), size=10.5, fill=BG1, line=BORDER,
        align=PP_ALIGN.CENTER)
    x += 2.38
box(s, L, 3.86, CW, 0.34, 'としポ基盤（よむすびパッケージ）／マルチテナント', size=12,
    color=WHITE, fill=PRIMARY, line=BORDER, align=PP_ALIGN.CENTER)
base_items = [('管理画面（よむすび管理画面）', 'メニュー・コンテンツ管理\nお知らせ／広報誌管理\n配信・承認管理'),
              ('配信基盤', 'プッシュ通知\n予約配信・エリア指定配信'),
              ('分析基盤（GCP／BigQuery）', '利用状況・KPIの可視化\nCSV出力・ダッシュボード')]
x = L
for t, d in base_items:
    box(s, x, 4.26, 3.06, 0.34, t, size=10.5, color=WHITE, fill=DARK, line=BORDER,
        align=PP_ALIGN.CENTER)
    box(s, x, 4.60, 3.06, 0.86, d.split('\n'), size=10, fill=BG1, line=BORDER,
        align=PP_ALIGN.CENTER)
    x += 3.17
label(s, L, 5.66, 1.70, 1.20, 'ご提案の\nポイント'.replace('\\n','\n').split('\n') if False else ['ご提案の','ポイント'])
textbox(s, 2.10, 5.66, 7.60, 1.20,
        ['決済（としポ）への到達速度を損なわない構成とし、ポイント機能は広島市モードに搭載しない',
         'メニュー・アイコン・遷移先はアプリ改修なしに管理画面から変更できる（Firebase Remote Config等を活用）',
         '防災は初年度リンク設置から開始し、R9年度以降に「避難所へGo!」同等機能の内製化へ段階的に拡張する'],
        size=11, bullets=True)

# ---------------------------------------------------------------- P5 4要件の全体整理（流用）
s = new_slide(prs); paste(s, 'scope')
replace_paragraph(s, {
 '4要件の全体整理｜実現する目的・背景と実現方針・スコープ外要素':
   '3．4要件の全体整理｜目的・背景／実現方針／今回のスコープ外',
 '4要件とも、としポ基盤（広島市モード）を土台に既存資産の活用と導線提供を基本としする。':
   '4要件とも、としポ基盤（広島市モード）を土台に、既存資産の活用と導線提供を基本方針とします。',
 '対象外とする内容（想定）': '今回のスコープ外（将来拡張として想定）',
 '既存電子申請サイトへの導線提供': '既存の「広島市オンライン手続ポータルサイト」への導線提供',
 '外部アプリ連携＋未インストール時はストア誘導': '外部アプリ連携＋未インストール時はアプリストアへ誘導',
})

# ---------------------------------------------------------------- P6 区切り
s = new_slide(prs)
box(s, 0.29, 2.90, 9.42, 0.05, '', fill=PRIMARY, line=None)
textbox(s, 0.29, 3.10, 9.42, 0.62, '4．機能概要と利用フロー', size=32,
        color=RGBColor(0,0,0), anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.29, 3.85, 9.42, 0.40,
        'A．広島市モード／B．電子申請／C．広報／D．防災', size=14, color=DARK)

# ---------------------------------------------------------------- P7 A/B/D フロー（流用）
s = new_slide(prs); paste(s, 'abd')
replace_paragraph(s, {
 'A.広島市モード／B.電子申請／D.防災｜機能概要と利用フロー':
   '4-1．A．広島市モード／B．電子申請／D．防災｜機能概要と利用フロー',
 '利用設定を起点にホーム（統合ポータル）から各サービスへ誘導し、配信・分析までを一体で運用する。':
   '利用設定を起点にホーム（統合ポータル）から各サービスへ誘導し、配信・分析までを一体で運用します。',
 'メニュー・アイコン・遷移先の変更はアプリ改修なしに管理画面から実施できる（項番89）':
   'メニュー・アイコン・遷移先の変更はアプリ改修なしに管理画面から実施できる（項番89）',
 'お知らせ配信は予約配信、承認フローを想定する': 'お知らせ配信は予約配信・承認フローに対応する',
 '※ フロー・役割分担は認識合わせ用のたたき台であり、ヒアリング結果を踏まえ要件定義で確定する。':
   '※ フロー・役割分担は現時点の想定であり、要件定義にて詳細を確定します。',
})

# ---------------------------------------------------------------- P8 C フロー（流用）
s = new_slide(prs); paste(s, 'c')
replace_paragraph(s, {
 'C.広報｜機能概要と利用フロー': '4-2．C．広報｜機能概要と利用フロー',
 '居住区の最新号をすぐ読める動線とし、入稿から公開・配信までを管理画面で完結させる。':
   '居住区の最新号をすぐ読める動線とし、入稿から公開・配信までを管理画面で完結させます。',
 '紙面はPDFのまま閲覧し、変換作業を行わない（項番58）': '紙面はPDFのまま閲覧し、変換作業を行わない（項番58）',
 '発行通知は統合ポータルの配信基盤と一体で運用する': '発行通知は統合ポータルの配信基盤と一体で運用する',
 '※ フロー・役割分担は認識合わせ用のたたき台であり、ヒアリング結果を踏まえ要件定義で確定する。':
   '※ フロー・役割分担は現時点の想定であり、要件定義にて詳細を確定します。',
})

# ---------------------------------------------------------------- P9 区切り
s = new_slide(prs)
box(s, 0.29, 2.90, 9.42, 0.05, '', fill=PRIMARY, line=None)
textbox(s, 0.29, 3.10, 9.42, 0.62, '5．画面イメージ（アプリ／管理画面）', size=32,
        color=RGBColor(0,0,0), anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.29, 3.85, 9.42, 0.40,
        '※ 詳細は別紙「アプリワイヤーフレーム」「管理画面ワイヤーフレーム」をご参照ください。',
        size=14, color=DARK)

# ---------------------------------------------------------------- P10 A ホーム
s = new_slide(prs); title(s, '5-1．A．広島市モード｜ホーム画面と表示設定')
lead(s, '広島市モード独自メニューと、としポ共通メニュー・お知らせを1画面に集約します。')
picture(s, f'{IMG}/app_home.png', 0.35, 1.25, h=4.90, caption='広島市モード ホーム')
picture(s, f'{IMG}/app_modeset.png', 2.65, 1.25, h=4.90, caption='広島市モード表示設定')
picture(s, f'{IMG}/app_rule.png', 5.05, 1.25, w=2.10, caption='ON／OFFの初期設定ルール')
textbox(s, 5.05, 2.60, 4.65, 3.60,
        ['ホーム上部に「広島市モード独自サービスメニュー」（電子申請／広報誌／休日保育予約／防災アプリ）を配置',
         '下部に「としポ共通サービスメニュー」（ゴミ出しカレンダー／クーポン／アンケート）を配置し、内容はとしポと同一',
         'お知らせは広島市向けのもののみを表示（エリア指定配信は新設）',
         '決済へは「としポ決済を開く」から即時遷移でき、到達速度を損なわない',
         'アカウントメニューに「広島市モード表示設定」を追加し、居住エリア・デフォルト表示（広報誌／お知らせ／ゴミ出しカレンダー）を設定',
         '郵便番号が未登録の場合は、初回にユーザへ確認のうえ設定'],
        size=10.5, bullets=True)
note(s, '出典：別紙「アプリワイヤーフレーム_draft1_20260824.pdf」（A・W）')

# ---------------------------------------------------------------- P11 A お知らせ・共通
s = new_slide(prs); title(s, '5-2．A．広島市モード｜お知らせと共通メニュー')
lead(s, 'エリア・カテゴリでの絞り込みと未読管理により、必要な情報に短い導線で到達できます。')
picture(s, f'{IMG}/app_oshirase.png', 0.35, 1.25, h=4.90, caption='お知らせ一覧（フィルタ）')
picture(s, f'{IMG}/app_gomi.png',     2.65, 1.25, h=4.90, caption='ゴミ出しカレンダー（既存）')
picture(s, f'{IMG}/app_account.png',  4.95, 1.25, h=4.90, caption='アカウントメニュー（既存）')
textbox(s, 7.25, 1.25, 2.50, 4.90,
        ['表示対象エリア（市全域／各区）で絞り込み。広島市モードのデフォルト指定に基づき初期表示',
         'カテゴリ（子育て／保育／防災／健康／イベント／その他）で絞り込み。初期値は全選択',
         '「未読のみ表示」フィルタを用意',
         'ゴミ出しカレンダーは広島市モードの表示設定に基づき地区をデフォルト表示',
         '未設定の場合は前回表示エリア、または共通デフォルトを適用'],
        size=10.5, bullets=True)
note(s, '出典：別紙「アプリワイヤーフレーム_draft1_20260824.pdf」（G・F・V）')

# ---------------------------------------------------------------- P12 B 電子申請
s = new_slide(prs); title(s, '5-3．B．電子申請｜画面イメージと連携方式')
lead(s, '既存の「広島市オンライン手続ポータルサイト」への導線を提供し、申請手続の入口をアプリに集約します。')
picture(s, f'{IMG}/app_eshinsei.png', 0.35, 1.25, h=5.00, caption='電子申請サイト（外部ブラウザ表示）')
label(s, 3.10, 1.25, 1.55, 0.80, ['連携方式'])
textbox(s, 4.80, 1.25, 4.95, 0.80,
        ['ホームの「電子申請」タップでブラウザを起動し、既存サイトを表示',
         '遷移先URL：https://apply.e-tumo.jp/city-hiroshima-u/offer/offerList_initDisplay'],
        size=10.5, bullets=True)
label(s, 3.10, 2.20, 1.55, 0.95, ['運用性'])
textbox(s, 4.80, 2.20, 4.95, 0.95,
        ['遷移先はリリース不要で変更可能とし、管理画面の設定項目またはFirebase Remote Configを活用',
         'URL変更・一時停止・案内文の差し替えを、アプリ改修を伴わずに実施可能'],
        size=10.5, bullets=True)
label(s, 3.10, 3.30, 1.55, 0.95, ['今回の','スコープ'])
textbox(s, 4.80, 3.30, 4.95, 0.95,
        ['既存電子申請サイトへの導線提供までを対象とする',
         '電子申請等APIとの連携によりアプリ内で手続を完結させる方式は、将来の拡張として位置づける'],
        size=10.5, bullets=True)
label(s, 3.10, 4.40, 1.55, 1.10, ['将来拡張','（参考）'])
textbox(s, 4.80, 4.40, 4.95, 1.10,
        ['電子申請等APIと連携し、ぴったりサービスの検索・申請をアプリ内で実施',
         'マイナンバーカードによる公的個人認証との連携',
         '請求書払い（納付書バーコード読取）などの機能拡張'],
        size=10.5, bullets=True)
note(s, '出典：別紙「アプリワイヤーフレーム_draft1_20260824.pdf」（B）／（修正版）20251126【広島市御中】よむすび機能説明資料')

# ---------------------------------------------------------------- P13 C 広報アプリ
s = new_slide(prs); title(s, '5-4．C．広報｜アプリ画面イメージ')
lead(s, '住所属性から居住区の最新号をデフォルト表示し、区・号の切り替えとバックナンバー閲覧に対応します。')
picture(s, f'{IMG}/app_kouhou.png',     0.35, 1.25, h=4.85, caption='広報（最新号）')
picture(s, f'{IMG}/app_kouhou_sel.png', 2.60, 1.25, h=4.85, caption='広報誌・号の選択')
picture(s, f'{IMG}/app_kouhou_dlg.png', 4.85, 1.25, h=4.85, caption='表示する区の選択')
textbox(s, 7.15, 1.25, 2.60, 4.85,
        ['住所属性に基づき居住区の広報紙をデフォルト表示（ユーザによる変更可）',
         '「表示する区を変更」から対象区を選択。チェックにより次回以降のデフォルトへ反映',
         '対象広報誌（市民と市政／ひろしま市議会だより／Hiroshima NOW 等）とバックナンバーを一覧から選択',
         '未読バッジにより新着を判別',
         '紙面はPDFのまま閲覧し、HTML化などの変換作業は行わない（項番58）',
         '取得は明示操作を起点とし、分割取得等で通信量・表示速度に配慮（項番59・60）'],
        size=10, bullets=True)
note(s, '出典：別紙「アプリワイヤーフレーム_draft1_20260824.pdf」（C）')

# ---------------------------------------------------------------- P14 C 管理画面1
s = new_slide(prs); title(s, '5-5．C．広報｜管理画面イメージ（広報誌管理・号数管理）')
lead(s, '広報誌の登録から号数・PDF・表紙の管理までを、既存のよむすび管理画面上で完結させます。')
picture(s, f'{IMG}/adm_list.png',   0.35, 1.22, w=4.60, caption='広報誌管理一覧')
picture(s, f'{IMG}/adm_gousuu.png', 5.10, 1.22, w=4.60, caption='号数管理／新しい号の登録・編集')
textbox(s, L, 4.85, CW, 1.95,
        ['広報誌ごとに発行元・対象区・最新号・ステータス（公開中／下書き）・最終更新日を一覧管理',
         '号数管理では、バックナンバーと最新号のPDF・表紙画像・ステータス・閲覧数を登録／管理',
         '新しい号の登録では、号数名・発行日・対象区（全区／各区）・PDFファイル（最大50MB）・表紙画像（推奨比率1：1.41）を設定',
         'プレビュー確認のうえ「号を保存して公開」。公開日を指定した予約公開にも対応',
         '広報誌の追加・削除、発行元による絞り込み、広報誌名での検索が可能'],
        size=10.5, bullets=True)
note(s, '出典：別紙「管理画面ワイヤーフレーム_draft1_20260824.pdf」')

# ---------------------------------------------------------------- P15 C 管理画面2
s = new_slide(prs); title(s, '5-6．C．広報｜管理画面イメージ（配信・承認管理）')
lead(s, '起案から承認・予約配信・プッシュ通知までを一連の流れで運用し、配信実績も同一画面で確認できます。')
picture(s, f'{IMG}/adm_haishin.png', 2.00, 1.22, w=6.00, caption='配信・承認管理')
textbox(s, L, 5.35, CW, 1.50,
        ['承認待ちの起案を一覧表示し、起案者・起案日・対象号を確認のうえPDFをプレビュー',
         'コメントを添えて「起案を却下・差し戻し」または「配信を承認・予約確定」',
         '配信時のプッシュ通知の同時送信可否、通知メッセージ（80文字まで）、配信予約日時を設定（空欄で即時配信）',
         '直近の配信完了履歴として、号数・配信日時・配信先区・対象者数・開封率を確認',
         '承認フローの権限（起案／承認）は、区役所・所管課などの運用体制に合わせて設定'],
        size=10.5, bullets=True)
note(s, '出典：別紙「管理画面ワイヤーフレーム_draft1_20260824.pdf」')

# ---------------------------------------------------------------- P16 D 防災
s = new_slide(prs); title(s, '5-7．D．防災｜「避難所へGo!」連携の画面イメージ')
lead(s, '既存の避難誘導アプリ「避難所へGo!」への入口を提供し、将来的な機能拡張の足がかりとします。')
picture(s, f'{IMG}/app_hinan_app.png',   0.35, 1.25, h=4.75, caption='インストール済み：アプリ起動')
picture(s, f'{IMG}/app_hinan_store.png', 2.60, 1.25, h=4.75, caption='未インストール：ストアへ遷移')
label(s, 4.90, 1.25, 1.40, 0.85, ['連携方式'])
textbox(s, 6.40, 1.25, 3.35, 0.85,
        ['インストール済み端末ではアプリを起動',
         '未インストール時はアプリストアの当該ページを表示'], size=10.5, bullets=True)
label(s, 4.90, 2.25, 1.40, 0.85, ['運用性'])
textbox(s, 6.40, 2.25, 3.35, 0.85,
        ['遷移先はリリース不要で変更可能とし、管理画面設定項目またはFirebase Remote Configを活用'],
        size=10.5, bullets=True)
label(s, 4.90, 3.25, 1.40, 1.05, ['今回の','スコープ'])
textbox(s, 6.40, 3.25, 3.35, 1.05,
        ['外部アプリ連携＋ストア誘導までを対象とする',
         'としポアプリ内での防災機能の実装、および「避難所へGo!」との機能連携は将来拡張として位置づける'],
        size=10.5, bullets=True)
label(s, 4.90, 4.45, 1.40, 1.05, ['次頁以降'])
textbox(s, 6.40, 4.45, 3.35, 1.05,
        ['R9年度以降の段階的な機能拡張について、ロードマップとしてご提案します（P18〜P20）'],
        size=10.5, bullets=True)
note(s, '出典：別紙「アプリワイヤーフレーム_draft1_20260824.pdf」（E）')

# ---------------------------------------------------------------- P17 区切り
s = new_slide(prs)
box(s, 0.29, 2.90, 9.42, 0.05, '', fill=PRIMARY, line=None)
textbox(s, 0.29, 3.10, 9.42, 0.62, '6．防災機能の拡張ロードマップ', size=32,
        color=RGBColor(0,0,0), anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.29, 3.85, 9.42, 0.40,
        '広島広域都市圏ポータルアプリ連携から、将来の機能追加まで', size=14, color=DARK)

# ---------------------------------------------------------------- P18 ロードマップ
s = new_slide(prs); title(s, '6-1．防災アプリ 機能拡張ロードマップ（案）')
lead(s, 'R8年度はリンク設置から着実に開始し、R9年度に同等機能の実装、R10年度以降に高度化を図ります。')
grid(s, L, 1.20, [1.30, 2.70, 2.70, 2.70],
     ['', 'R8年度', 'R9年度', 'R10年度以降'],
     [['実装機能\n（予定）',
       'ポータルアプリから「避難所へGo!」へのリンク設置',
       '「避難所へGo!」と同等機能の実装',
       '浸水被害ARシミュレーション機能\n避難所受付システム\n独自防災システムとのAPI連携（双方向データ連携）'],
      ['検討事項',
       '①「避難所へGo!」と同等機能の実装\n②浸水被害ARシミュレーション\n③避難所受付システム',
       '①浸水被害ARシミュレーション\n②避難所受付システム\n③貴市防災システムとのAPI連携\n④関連市町・周辺市町への拡大対応',
       '－'],
      ['ご提案内容',
       '上記①②③の実装に関する概算お見積もりの提出（期日：2026年9月末）',
       '上記①②③の実装に関する概算お見積もりの提出',
       '－'],
      ['補足',
       '②③は詳細要件が未定のため参考試算\n関連市町（廿日市市・江田島市・熊野町）への事前説明が必要',
       '避難所受付システムや防災システムとのAPI連携は、貴市の防災計画等を踏まえた全体的な整理が必要',
       '－']],
     heights=[1.35, 1.35, 1.05, 1.35], size=10)
note(s, '※ 同等機能の詳細は別紙「20260824_現状機能概要まとめ.xlsx」で整理しています。「避難所へGo!」は廿日市市・江田島市・熊野町でも利用中です。')

# ---------------------------------------------------------------- P19 同等機能の範囲
s = new_slide(prs); title(s, '6-2．R9年度｜「避難所へGo!」同等機能の実装範囲')
lead(s, '現行アプリの機能を項番単位で棚卸しし、ミニアプリとして取り込む範囲を整理しています。')
grid(s, L, 1.20, [2.10, 5.55, 1.75],
     ['機能分類', '主な機能', '項番'],
     [['全般', '多言語表示（日本語／英語／中国語（繁体・簡体）／韓国語）。将来的な対応言語の拡充も可能', '1'],
      ['アプリ初回起動', '位置情報・通知等の権限取得、アプリ非起動時のGPS ON／OFF設定', '2〜9'],
      ['地域選択', '小学校区単位での防災情報受信地域の設定。現在地に加え事前登録地点を3か所程度選択', '10〜13'],
      ['トップメニュー', 'ロゴ表示、現在地の天気予報・最高最低気温、ボトムメニュー、リンク集', '14〜16'],
      ['危険度の確認', 'マップ基図切替（Googleマップ／地理院地図）、ハザード情報のレイヤ表示（土砂・河川・津波・高潮・洪水・雨水・ため池）、危険性判定、オフライン地図、ARカメラ・避難コンパスによる避難誘導', '17〜23'],
      ['最寄りの避難所へ', 'Lアラート連携による開設避難所への最短ルート検索、経由地検索・ルート再検索', '24〜25'],
      ['安否登録・確認', 'web171・Googleパーソンファインダー・J-anpiへのリンク、家族へのメール送信', '26〜31'],
      ['防災ハンドブック', 'アプリ上での防災ハンドブックの閲覧', '32〜33'],
      ['避難所検索', '現在地・地図上の任意地点・全国の地域リストからの検索、避難所詳細（名称・住所・標高・災害対応種別・収容人数・備蓄）', '34〜35'],
      ['防災情報・お知らせ', 'Lアラート連携の防災緊急情報、履歴100件程度の保存・未読表示、自治体情報・気象情報の提供', '36〜37']],
     heights=[0.42,0.42,0.50,0.42,0.78,0.50,0.42,0.34,0.62,0.50], size=9.5)
note(s, '※ 地域選択の対象は広島市8区（150区分）・廿日市市（28）・江田島市（25）・熊野町（1）の計204区分。項番は別紙「20260824_現状機能概要まとめ.xlsx」に対応します。')

# ---------------------------------------------------------------- P20 防災×としポ
s = new_slide(prs); title(s, '6-3．防災×としポ 企画案')
lead(s, '防災アプリの普及率向上と平時からの防災行動の定着を、としポの付与施策と組み合わせて実現します。')
grid(s, L, 1.20, [2.30, 4.20, 2.90],
     ['項目', '内容', '補足'],
     [['防災アプリの\n普及率向上',
       '下記の行動に対してとしポを付与\n「避難所へGo!」のインストール（実施中）\n「避難所へGo!」を開く\n広島市防災情報メールの登録\n災害ボランティア事前登録',
       'としポとは別に「防災ポイント」（1p＝1円ではないポイント）を発行可能\n防災グッズ等が当たる抽選への応募といった施策も実施可能'],
      ['避難所チェックラリー',
       'マップ上に表示された避難所等の半径●m以内に入るとチェックインが可能（1日1回など）\nチェックインで1日1回ポイントを付与',
       'チェックイン＋防災に関する動画閲覧でのポイント付与も可能'],
      ['防災スタンプラリー',
       '避難所巡りでスタンプを獲得し、達成者にとしポを付与\n広域都市圏を対象としたスタンプラリーも可能',
       '例）広島市中心部：広島市役所前避難所、平和記念公園津波避難ビル、比治山公園'],
      ['防災クイズ',
       '3択クイズに正解するととしポをプレゼント',
       '例）広島市で特に多い災害はどれですか？（大規模雪害／土砂災害／竜巻）\n例）大雨危険警報はどの警戒レベルに対応しますか？（警戒レベル4／3／2）']],
     heights=[1.50, 1.05, 1.00, 0.95], size=10)
note(s, '★ 貯まった防災ポイントで防災グッズ等が当たる抽選への応募、一定以上の獲得者を「防災マスター」として表彰する等の運用も可能です。')

# ---------------------------------------------------------------- P21 管理画面と運用
s = new_slide(prs); title(s, '7-1．管理画面と運用フロー')
lead(s, 'アプリ改修を伴わない管理画面運用により、市の運用サイクルに合わせた情報発信を実現します。')
grid(s, L, 1.20, [2.40, 3.60, 3.40],
     ['管理メニュー', '主な操作', '運用イメージ'],
     [['メニュー・\nコンテンツ管理', 'メニュー名称／アイコン／並び順／遷移先の登録・変更',
       'アプリ改修・ストア申請なしで反映\n（項番89）'],
      ['お知らせ管理', 'お知らせの起案・カテゴリ設定・対象エリア指定（市全域／各区）',
       '所管課が起案し、広報担当が承認'],
      ['広報誌管理', '広報誌・号数の登録、PDF／表紙画像の登録、公開ステータス管理',
       '入稿から公開までを1画面で完結'],
      ['配信・承認管理', '起案の承認／却下・差し戻し、予約配信日時の設定、プッシュ通知の同時送信',
       '承認フローにより誤配信を防止'],
      ['通知管理', 'プッシュ通知の作成・予約配信・配信結果の確認',
       '発行案内の通知と広報誌公開を連動'],
      ['利用状況の確認', '配信対象者数・開封率・閲覧数の確認、CSV出力',
       '効果検証と次回施策の改善に活用']],
     heights=[0.62,0.55,0.55,0.62,0.48,0.48], size=10)
note(s, '※ 権限（起案／承認／閲覧）は操作ID管理により、所管課・区役所等の運用体制に合わせて設定します。')

# ---------------------------------------------------------------- P22 分析
s = new_slide(prs); title(s, '7-2．分析・効果測定')
lead(s, 'GCP／BigQueryを活用し、利用状況の可視化とデータに基づく施策改善をご支援します。')
label(s, L, 1.25, 1.70, 1.05, ['分析基盤'])
textbox(s, 2.10, 1.25, 7.60, 1.05,
        ['弊社システムはGCPのサーバー環境で稼働し、BigQueryを活用したデータの可視化が可能',
         'ダッシュボードは事業特性に応じてカスタマイズ可能。管理画面から取引データのCSV出力にも対応'],
        size=11, bullets=True)
label(s, L, 2.45, 1.70, 1.05, ['分析観点'])
textbox(s, 2.10, 2.45, 7.60, 1.05,
        ['ファネル分析により、登録者数の伸び悩みや同意画面での離脱等の原因を究明',
         'アプリ・サーバーの各種データを取り込み、施策実施前後の変化を定量的に確認'],
        size=11, bullets=True)
textbox(s, L, 3.70, CW, 0.36, 'KPI（案）', size=12, color=DARK, bold=True)
grid(s, L, 4.10, [2.30, 3.70, 3.40],
     ['区分', '指標', '取得方法'],
     [['A．広島市モード', 'インストール数／広島市モードON率／ホーム表示回数', '弊社システムDB／BigQuery'],
      ['B．電子申請', '電子申請メニューの遷移数・遷移率', 'アプリイベントログ'],
      ['C．広報', '広報紙の閲覧数／通知開封率／区別の閲覧傾向', '管理画面（閲覧数・開封率）'],
      ['D．防災', '「避難所へGo!」への遷移数／ストア遷移数', 'アプリイベントログ']],
     heights=[0.55,0.42,0.55,0.42], size=10)
note(s, '※ 指標および取得方法は、要件定義にて最終確定します。')

# ---------------------------------------------------------------- P23 基盤・セキュリティ
s = new_slide(prs); title(s, '7-3．システム基盤・セキュリティ')
lead(s, '既存パッケージの実績ある基盤を活用し、可用性・セキュリティ・運用性を確保します。')
grid(s, L, 1.20, [2.30, 7.10],
     ['区分', '内容'],
     [['稼働環境',
       'インターネット上のクラウド（GCP）を利用し、安定的な常時運用が可能で耐災性の高い構成とする\nデータセンターは、広島市で災害が発生しても影響が生じない地域のものとする'],
      ['認証・第三者評価',
       'プライバシーマーク、情報セキュリティマネジメントシステム（ISMS／ISO27001）の認証を取得\nデジタル庁「デジタル地方創生サービスカタログ」に「よむすび」「健康ポイント」が掲載'],
      ['アプリの安全性',
       'PDF閲覧時のセキュリティに配慮（外部ファイルに起因するリスク対策）\n対象OSは、リリース時点でサポートされる版数に合わせて設定'],
      ['運用性',
       '電子申請・防災アプリの遷移先は、Firebase Remote Config等の標準機能を活用しリリース不要で変更可能\nメニュー・コンテンツは管理画面から変更でき、アプリ改修・ストア申請を伴わない'],
      ['防災同等機能\n実装時の前提',
       '10万人が同時にアクセスした場合も支障なく運用できる容量と性能を確保\nポップアップ通知は、正常な通信下において10万人に対して10秒程度で全利用者へ通知\nLアラートの情報伝達者申請、App Store／Google Playへの登録申請の補助を実施']],
     heights=[0.85,0.85,0.85,0.85,1.10], size=10)
note(s, '※ 「防災同等機能実装時の前提」は、現行「避難所へGo!」の仕様書に基づく要件です（R9年度以降のご提案範囲）。')

# ---------------------------------------------------------------- P24 スケジュール
s = new_slide(prs); title(s, '8-1．実施スケジュール（案）')
lead(s, 'R8年度は4要件のリリースを最優先とし、並行して防災機能拡張の概算お見積もりを提出します。')
months = ['2026/9','10月','11月','12月','2027/1','2月','3月','R9年度以降']
cx = 2.20; cw = 0.94
box(s, L, 1.25, 1.90, 0.34, '工程', size=11, color=WHITE, fill=PRIMARY, line=BORDER,
    align=PP_ALIGN.CENTER)
for m in months:
    box(s, cx, 1.25, cw, 0.34, m, size=9.5, color=WHITE, fill=PRIMARY, line=BORDER,
        align=PP_ALIGN.CENTER)
    cx += cw
tasks = [('要求仕様の確定', 0, 1), ('要件定義・基本設計', 1, 2), ('詳細設計・開発', 2, 3),
         ('テスト（結合・総合）', 5, 1), ('ストア申請・リリース', 6, 1),
         ('防災：概算お見積もり提出', 0, 1), ('防災：同等機能の実装', 7, 1)]
ty = 1.59
for name, st, ln in tasks:
    box(s, L, ty, 1.90, 0.46, name, size=10, color=DARK, fill=BG2, line=BORDER,
        align=PP_ALIGN.CENTER)
    for i in range(len(months)):
        box(s, 2.20 + cw * i, ty, cw, 0.46, '', fill=BG1 if i % 2 else WHITE, line=BORDER)
    box(s, 2.24 + cw * st, ty + 0.11, cw * ln - 0.08, 0.24, '', fill=PRIMARY, line=None)
    ty += 0.46
textbox(s, L, ty + 0.15, CW, 1.30,
        ['防災機能拡張（同等機能の実装／浸水被害ARシミュレーション／避難所受付システム）の概算お見積もりは、2026年9月末までに提出します',
         '関連市町（廿日市市・江田島市・熊野町）への事前説明が必要となるため、実施時期は貴市と調整のうえ決定します',
         '「休日保育予約」は別資料にて整理のうえ、スケジュールを別途ご提示します'],
        size=10.5, bullets=True)
note(s, '※ 期間は現時点の目安であり、要求仕様の確定内容および貴市の調達手続に応じて調整します。')

# ---------------------------------------------------------------- P25 今後の進め方
s = new_slide(prs); title(s, '8-2．今後の進め方とご確認事項')
lead(s, '以下の事項についてご確認・ご指示をいただき、要求仕様を確定させていただきたく存じます。')
grid(s, L, 1.20, [0.75, 2.55, 6.10],
     ['No', '区分', 'ご確認事項'],
     [['1', 'A．広島市モード', '広島市モードのON／OFF初期設定ルール（郵便番号登録済・広島市内＝ON／市外＝OFF／未登録＝ユーザに確認）でよいか'],
      ['2', 'A．広島市モード', 'お知らせのエリア区分（市全域／各区）とカテゴリ体系（子育て／保育／防災／健康／イベント／その他）の妥当性'],
      ['3', 'B．電子申請', '遷移先URLおよび遷移方式（外部ブラウザ起動）でよいか。認証連携の要否'],
      ['4', 'C．広報', '対象とする広報紙（市民と市政／ひろしま市議会だより／区報／Hiroshima NOW 等）、発行元・対象区、バックナンバーの保持期間'],
      ['5', 'C．広報', '入稿・公開・配信承認フローの運用主体と権限（起案／承認）の設定'],
      ['6', 'D．防災', '「避難所へGo!」への遷移方式（アプリ起動／ストア誘導）でよいか'],
      ['7', 'D．防災', 'R9年度以降の機能拡張の優先順位、および関連市町への事前説明の進め方'],
      ['8', '共通', '「休日保育予約」との連携範囲（別資料にて整理）']],
     heights=[0.62,0.62,0.48,0.62,0.48,0.40,0.48,0.40], size=10)
note(s, '※ ご回答内容を踏まえ、要求仕様書として整理のうえ改めてご提示いたします。')

# ---------------------------------------------------------------- P26 別紙一覧
s = new_slide(prs); title(s, '別紙一覧')
lead(s, '本資料は、以下の別紙とあわせてご確認ください。')
grid(s, L, 1.25, [3.40, 6.00],
     ['資料名', '内容'],
     [['アプリワイヤーフレーム\n_draft1_20260824.pdf',
       '広島市モードホーム、電子申請、広報誌、防災アプリ連携、お知らせ、広島市モード表示設定のアプリUIイメージおよび要求仕様の注記'],
      ['管理画面ワイヤーフレーム\n_draft1_20260824.pdf',
       '広報誌管理一覧、号数管理（新しい号の登録・編集）、配信・承認管理の管理画面UIイメージ'],
      ['20260824_現状機能概要まとめ.xlsx',
       '「避難所へGo!」の現行機能を項番単位で整理した一覧（仕様書該当内容・現行アプリ挙動確認結果・重要性評価）'],
      ['（参考）追記_20260814_避難所へGo!\n機能拡張ロードマップ(案).pptx',
       '防災アプリの機能拡張ロードマップおよび防災×としポ企画案'],
      ['（参考）（修正版）20251126【広島市御中】\nよむすび機能説明資料.pptx',
       'スーパーアプリのコンセプト、ミニアプリ（電子申請・住民との連絡機能・各種防災機能）、分析機能']],
     heights=[0.90,0.72,0.80,0.80,0.80], size=10)
textbox(s, L, 6.30, CW, 0.60,
        ['本ご提案に関するお問い合わせ先：フェリカポケットマーケティング株式会社'],
        size=10.5, color=DARK)

# ---------------------------------------------------------------- ページ番号
for i, sl in enumerate(prs.slides):
    if i == 0: continue
    pagenum(sl, i + 1)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print('saved', OUT, len(prs.slides.__iter__.__self__._sldIdLst), 'slides')
