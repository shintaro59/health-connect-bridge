# -*- coding: utf-8 -*-
"""広島市 防災アプリ提案書 共通スタイル"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FONT   = 'Meiryo UI'
PRIMARY= RGBColor(0x00,0x81,0xAF)
DARK   = RGBColor(0x00,0x58,0x7A)
BORDER = RGBColor(0xD8,0xE2,0xE6)
BG1    = RGBColor(0xF7,0xFA,0xFB)
BG2    = RGBColor(0xF0,0xF5,0xF7)
TEXT   = RGBColor(0x26,0x26,0x26)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GRAY   = RGBColor(0x77,0x88,0x8E)
BLACK  = RGBColor(0x00,0x00,0x00)
L, CW  = 0.30, 9.40

def _rpr_font(run, size, color, bold=False):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', FONT)

def _bullet(par, on=True, size=10.5, color=PRIMARY):
    pPr = par._p.get_or_add_pPr()
    for t in ('a:buNone','a:buChar','a:buClr','a:buSzPct','a:buFont','a:lnSpc','a:spcAft','a:spcBef'):
        for e in pPr.findall(qn(t)): pPr.remove(e)
    def mk(tag, **kw): return pPr.makeelement(qn(tag), kw)
    ln = mk('a:lnSpc'); ln.append(ln.makeelement(qn('a:spcPct'), {'val':'112000'})); pPr.append(ln)
    sb = mk('a:spcBef'); sb.append(sb.makeelement(qn('a:spcPts'), {'val':'0'})); pPr.append(sb)
    sa = mk('a:spcAft'); sa.append(sa.makeelement(qn('a:spcPts'), {'val':'300'})); pPr.append(sa)
    if on:
        ind = int(Pt(size).emu * 1.30); mar = ind + int(Pt(size).emu*0.22)
        pPr.set('marL', str(mar)); pPr.set('indent', str(-ind))
        bc = mk('a:buClr'); bc.append(bc.makeelement(qn('a:srgbClr'), {'val':'%02X%02X%02X'%(color[0],color[1],color[2])})); pPr.append(bc)
        pPr.append(mk('a:buSzPct', val='55000'))
        pPr.append(mk('a:buFont', typeface='Arial'))
        pPr.append(mk('a:buChar', char='●'))
    else:
        pPr.set('marL','0'); pPr.set('indent','0'); pPr.append(mk('a:buNone'))

def _fill_tf(tf, lines, size, color, bold, align, anchor, bullets):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        opt = {}
        if isinstance(item, tuple): item, opt = item
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = opt.get('align', align)
        sz = opt.get('size', size)
        _bullet(par, on=opt.get('bullets', bullets), size=sz, color=opt.get('bucolor', PRIMARY))
        r = par.add_run(); r.text = item
        _rpr_font(r, sz, opt.get('color', color), opt.get('bold', bold))

def textbox(slide, x, y, w, h, lines, size=11, color=TEXT, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, bullets=False):
    if isinstance(lines, str): lines = [lines]
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.margin_top = tf.margin_bottom = Emu(18288)
    _fill_tf(tf, lines, size, color, bold, align, anchor, bullets)
    return tb

def box(slide, x, y, w, h, lines, size=10.5, color=TEXT, fill=None, line=BORDER,
        bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, bullets=False,
        shape=MSO_SHAPE.RECTANGLE):
    if isinstance(lines, str): lines = [lines]
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(0.5)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = Emu(64008)
    tf.margin_top = tf.margin_bottom = Emu(27432)
    _fill_tf(tf, lines, size, color, bold, align, anchor, bullets)
    return sp

def arrow(slide, x, y, w=0.20, h=0.22, color=PRIMARY, shape=MSO_SHAPE.RIGHT_ARROW):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    return sp

def title(slide, text):
    return textbox(slide, L, 0.15, 9.70, 0.38, text, size=20, color=BLACK,
                   anchor=MSO_ANCHOR.MIDDLE)

def lead(slide, text, y=0.63, size=12):
    if isinstance(text, str): text = [text]
    return textbox(slide, L, y, CW, 0.42 if len(text) == 1 else 0.68, text,
                   size=size, color=TEXT)

def note(slide, text, y=6.92, size=9):
    return textbox(slide, L, y, CW, 0.30, text, size=size, color=GRAY)

def secthead(slide, x, y, w, text, size=11, h=0.34, fill=None):
    return box(slide, x, y, w, h, text, size=size, color=WHITE,
               fill=fill or PRIMARY, line=BORDER, align=PP_ALIGN.CENTER)

def label(slide, x, y, w, h, lines, size=11):
    if isinstance(lines, str): lines = [lines]
    return box(slide, x, y, w, h, lines, size=size, color=WHITE, fill=DARK,
               line=BORDER, align=PP_ALIGN.CENTER)

def grid(slide, x, y, widths, header, rows, hh=0.40, heights=None, size=10,
         first_col_label=True):
    cx = x
    for w, t in zip(widths, header):
        secthead(slide, cx, y, w, t, size=size + 0.5, h=hh); cx += w
    cy = y + hh
    for ri, row in enumerate(rows):
        rh = (heights[ri] if heights else 0.55)
        cx = x
        for ci, (w, cell) in enumerate(zip(widths, row)):
            first = (ci == 0 and first_col_label)
            lines = [l for l in str(cell).split('\n') if l.strip() != '']
            multi = len(lines) > 1
            box(slide, cx, cy, w, rh, lines if lines else '',
                size=size,
                color=DARK if first else TEXT,
                fill=BG2 if first else (BG1 if ri % 2 else WHITE),
                bold=first,
                align=PP_ALIGN.CENTER if first else PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE if (first or not multi) else MSO_ANCHOR.TOP,
                bullets=(not first) and multi)
            cx += w
        cy += rh
    return cy

def clear(slide):
    tree = slide.shapes._spTree
    for sh in list(slide.shapes):
        tree.remove(sh._element)

def pagenum(slide, n):
    textbox(slide, 9.05, 7.05, 0.70, 0.26, str(n), size=9, color=GRAY,
            align=PP_ALIGN.RIGHT)

def para_text(par):
    return ''.join(r.text for r in par.runs)

def drop_empty_paragraphs(shape):
    """本文中の空段落（空の箇条書き）を削除"""
    if not shape.has_text_frame: return
    tf = shape.text_frame
    ps = list(tf.paragraphs)
    for par in ps[1:] if len(ps) > 1 else []:
        if para_text(par).strip() == '':
            par._p.getparent().remove(par._p)
    if len(tf.paragraphs) > 1 and para_text(tf.paragraphs[0]).strip() == '':
        tf.paragraphs[0]._p.getparent().remove(tf.paragraphs[0]._p)

def replace_paragraph(slide, mapping):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for par in sh.text_frame.paragraphs:
            t = para_text(par)
            if t in mapping:
                runs = par.runs
                if not runs: continue
                runs[0].text = mapping[t]
                for r in runs[1:]:
                    r._r.getparent().remove(r._r)

def find_shape(slide, text):
    for sh in slide.shapes:
        if sh.has_text_frame and text in sh.text_frame.text:
            return sh
    return None

def delete_slide(prs, idx):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    prs.part.drop_rel(ids[idx].get(qn('r:id')))
    lst.remove(ids[idx])

def reorder(prs, order):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for e in ids: lst.remove(e)
    for i in order: lst.append(ids[i])
