# -*- coding: utf-8 -*-
import sys, os
sys.path.insert(0, '/tmp/w')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from deckkit import *
from deckkit import _rpr_font, _bullet

SRC = '/root/.claude/uploads/33b70c79-5330-5515-8cdd-be08ae2c87f2/8453af26-20260827_______________.pptx'
OUT = '/tmp/w/out2/20260827_広島市危機管理対策課御中_防災アプリ実装・機能拡張のご提案.pptx'

prs = Presentation(SRC)
S = list(prs.slides)

# ================================================================ P1 表紙
s = S[0]
box(s, 0.30, 3.55, 2.20, 0.05, '', fill=PRIMARY, line=None)
textbox(s, 0.28, 3.75, 8.50, 0.45,
        '〜 R9年度「避難所へGo!」同等機能の実装と、その先の機能拡張 〜',
        size=15, color=DARK)

# ================================================================ P2 目的とゴール（表現の統一）
s = S[1]
lead2 = find_shape(s, '認識合わせ')
if lead2 is not None: lead2.width = Inches(9.40)
replace_paragraph(s, {
 '防災アプリ実装・機能拡張について、認識合わせのためたたき台を提示する。':
   '防災アプリの実装・機能拡張について、現時点の想定をご提示し、認識合わせをさせていただく。',
 'R8年度はポータルアプリから「避難所へGo!」へのリンク設置':
   'R8年度はポータルアプリから「避難所へGo!」へのリンク設置。R9年度に同等機能の実装を予定',
 '※ 本資料は要求仕様確定前のたたき台であり、記載内容は今後のヒアリングで精査・確定する。':
   '※ 本資料は要求仕様確定前の想定であり、記載内容は今後のヒアリングを踏まえ精査・確定する。',
})

# ================================================================ 表のスタイル統一
def style_table(shape, header_size=10, body_size=9, first_col_label=True):
    tbl = shape.table
    tbl.first_row = False; tbl.horz_banding = False; tbl.first_col = False
    for ri, row in enumerate(tbl.rows):
        for ci, cell in enumerate(row.cells):
            first = first_col_label and ci == 0
            cell.fill.solid()
            if ri == 0:   cell.fill.fore_color.rgb = PRIMARY
            elif first:   cell.fill.fore_color.rgb = BG2
            else:         cell.fill.fore_color.rgb = BG1 if ri % 2 else WHITE
            cell.margin_left = cell.margin_right = Emu(64008)
            cell.margin_top = cell.margin_bottom = Emu(27432)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for par in cell.text_frame.paragraphs:
                sz  = header_size if ri == 0 else body_size
                col = WHITE if ri == 0 else (DARK if first else TEXT)
                bd  = (ri == 0) or first
                par.alignment = PP_ALIGN.CENTER if (ri == 0 or first) else PP_ALIGN.LEFT
                for r in par.runs: _rpr_font(r, sz, col, bd)

# ---- P3 ロードマップ
s = S[2]
sh = find_shape(s, '18')
if sh is not None and sh.text_frame.text.strip() == '18':
    sh._element.getparent().remove(sh._element)          # 前資料の残存ページ番号
for sh in s.shapes:
    if sh.has_table:
        style_table(sh, header_size=11, body_size=9.5)
        rows = list(sh.table.rows)
        for r, h in zip(rows, [0.42, 1.15, 1.60, 0.85, 1.30]): r.height = Inches(h)
        sh.top = Inches(1.10)
sub = find_shape(s, '広島広域都市圏ポータルアプリ連携から')
if sub is not None:
    for par in sub.text_frame.paragraphs:
        for r in par.runs: _rpr_font(r, 12, TEXT, False)
note(s, '※ 同等機能の詳細は別紙「20260824_現状機能概要まとめ.xlsx」にて整理しています。'
        '「避難所へGo!」は廿日市市・江田島市・熊野町でも利用中です。')

# ---- P4 防災×としポ
s = S[3]
sh = find_shape(s, '20')
if sh is not None and sh.text_frame.text.strip() == '20':
    sh._element.getparent().remove(sh._element)
for sh in s.shapes:
    if sh.has_table:
        style_table(sh, header_size=11, body_size=10)
        sh.top = Inches(1.15)                              # リード文との重なりを解消
        rows = sh.table.rows
        rows[0].height = Inches(0.42)
        for r in list(rows)[1:]: r.height = Inches(1.12)
star = find_shape(s, '★としポとは別の')
if star is not None:
    star.top = Inches(6.15)
    for par in star.text_frame.paragraphs:
        for r in par.runs: _rpr_font(r, 9.5, GRAY, False)

# ================================================================ P5 スコープと前提
s = S[4]
for sh in s.shapes: drop_empty_paragraphs(sh)             # 空の箇条書きを除去
zen = None
for sh in s.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith('防災アプリ、管理システムの要求仕様確定'):
        zen = sh
if zen is not None:
    zen.text_frame.word_wrap = True
    p = zen.text_frame.paragraphs[-1]
    for extra in ['「避難所へGo!」の権利関係および現行ベンダーとの調整方針の整理',
                  '関連市町（廿日市市・江田島市・熊野町）への事前説明']:
        np = zen.text_frame.add_paragraph()
        _bullet(np, on=True, size=11)
        r = np.add_run(); r.text = extra; _rpr_font(r, 11, TEXT, False)

# ================================================================ P6 仕様書概要（行の高さ・位置ずれ）
s = S[5]
FIX = {'個別機能要件': (3.09, 0.95), '職員研修': (4.06, 0.55),
       'その他関連業務': (4.63, 0.55), 'プロジェクト管理': (5.20, 0.55)}
for sh in s.shapes:
    if not sh.has_text_frame: continue
    t = sh.text_frame.text.strip()
    if t in FIX and sh.width < Inches(2.0):
        top, h = FIX[t]
        sh.top = Inches(top); sh.height = Inches(h)
    if t.startswith('防災情報等の受信地域'):
        sh.top = Inches(3.09); sh.height = Inches(0.95)
    if t.startswith('操作マニュアルをもとに'):
        sh.top = Inches(4.06); sh.height = Inches(0.55)
    if t.startswith('運用経費、バグ対応'):
        sh.top = Inches(4.63); sh.height = Inches(0.55)
    if t.startswith('プロジェクト管理、課題管理'):
        sh.top = Inches(5.20); sh.height = Inches(0.55)
for sh in list(s.shapes):
    if sh.top is not None and sh.top >= Inches(0.80):
        sh.top = sh.top + Inches(0.22)
lead(s, '現行「避難所へGo!」の仕様書に定められた要件の全体像。R9年度はこの範囲と同等の機能を実装対象とする。')

# ================================================================ P8 利用フロー（リード文の追加）
s = S[7]
lead(s, '平時の設定から、災害発生時の情報取得・通知・避難までを一連のフローとして想定している。')

# ================================================================ P9 フローのバリエーション（空表を記入）
s = S[8]
clear(s)
title(s, 'フローのバリエーション')
lead(s, '通常フローに加え、想定される分岐・例外のパターンを整理する。')
grid(s, L, 1.30, [1.60, 4.30, 3.50],
     ['業務フェーズ', '想定されるバリエーション（分岐・例外）', '備考'],
     [['初回起動・設定', '位置情報／通知の権限を許可しない場合。地域を手動選択して利用',
       '一部機能に制限。地域選択のみでも防災情報の受信は可能（項番10・11）'],
      ['平時', '自宅・勤務先・実家など、現在地以外の地点でも通知を受け取りたい',
       '事前登録地点を3か所程度設定し、現在地とあわせて受信（項番4・13）'],
      ['警戒・発災時', '通信が不安定または圏外で、地図・避難所情報を取得できない',
       '表示済みの地理院地図をキャッシュし、オフラインで表示（項番18）'],
      ['警戒・発災時', '配信エリア外にいて通知を受け取れなかった利用者が、後からエリアに進入',
       '一定期間内であれば遡って受信（GPS連動Push・期間指定配信）'],
      ['避難行動', 'ルート上に危険箇所があり通過できない',
       '経由地検索・ルート再検索により別ルートを案内（項番25）'],
      ['避難行動', 'カメラの権限を許可しない、またはAR非対応端末である',
       'ARカメラは利用不可。避難コンパスの表示で代替（項番23）'],
      ['安否確認', '家族が別の場所にいる／端末を持っていない',
       'web171・Googleパーソンファインダー・J-anpiへのリンク、家族へのメール送信（項番26〜31）']],
     heights=[0.72]*7, size=9.5)
note(s, '※ 項番は別紙「20260824_現状機能概要まとめ.xlsx」に対応します。バリエーションはヒアリングを踏まえて追加・確定します。')

# ================================================================ P10 ヒアリング論点（空の箇条書き・配色）
s = S[9]
for sh in s.shapes: drop_empty_paragraphs(sh)
for sh in s.shapes:
    if sh.has_text_frame and 'R10/3/xx' in sh.text_frame.text:
        sh.fill.solid(); sh.fill.fore_color.rgb = DARK
        sh.line.color.rgb = BORDER; sh.line.width = Pt(0.5)
        for par in sh.text_frame.paragraphs:
            for r in par.runs: _rpr_font(r, 11, WHITE, True)
ronten = find_shape(s, '想定している利用フロー以外の利用フロー')
if ronten is not None:
    tf = ronten.text_frame
    for extra in ['現行アプリのデータ（避難所等情報、小学校区、ハザード情報）の提供方法と更新頻度',
                  '関連市町（廿日市市・江田島市・熊野町）への説明・調整の進め方']:
        np = tf.add_paragraph(); _bullet(np, on=True, size=11)
        r = np.add_run(); r.text = extra; _rpr_font(r, 11, TEXT, False)

# ================================================================ P11 浸水被害ARシミュレーション（新規作成）
s = S[10]
clear(s)
title(s, '機能拡張イメージ①｜浸水被害ARシミュレーション')
lead(s, 'スマートフォンのカメラで日常の風景を映すだけで、その場所が「どこまで浸水するのか」を実寸大で体感できる機能。')
secthead(s, L, 1.15, 9.40, '利用イメージ', size=11, h=0.32)
steps = [('① カメラを向ける',   'アプリの「浸水ARを見る」から、\n自宅前や通学路などにカメラを向ける'),
         ('② 浸水面を重ねて表示', '地面・床を認識し、指定した深さの\n水面を実寸大で合成表示する'),
         ('③ 深さを変えて共有',   '浸水深を切り替えて比較。静止画・\n動画を保存して家族や地域で共有')]
x = L
for i, (h, d) in enumerate(steps):
    box(s, x, 1.55, 2.90, 1.05,
        [(h, {'bold': True, 'color': DARK, 'size': 11, 'align': PP_ALIGN.CENTER})]
        + [(l, {'size': 9.5, 'align': PP_ALIGN.CENTER}) for l in d.split('\n')],
        fill=BG1, line=BORDER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 2: arrow(s, x + 2.94, 1.96, 0.20, 0.22)
    x += 3.14

secthead(s, L, 2.78, 4.45, '浸水深の目安（表示を切り替えて比較）', size=10.5, h=0.32)
depths = [('0.5m',    RGBColor(0xCF,0xE7,0xF0), DARK,  '大人のひざ程度。歩行が難しくなり、自動車も走行が困難になる'),
          ('1.0m',    RGBColor(0x7F,0xC0,0xD7), DARK,  '大人の腰程度。車のドアが開かなくなり、徒歩での避難は危険'),
          ('2.0m',    RGBColor(0x00,0x81,0xAF), WHITE, '1階の軒下まで浸水。1階は水没し、2階への垂直避難が必要'),
          ('3.0m以上', RGBColor(0x00,0x58,0x7A), WHITE, '2階の床上まで浸水。垂直避難では不十分で、事前の立退き避難が必要')]
y = 3.10
for d, c, fg, t in depths:
    box(s, L, y, 1.00, 0.78, d, size=11, color=fg, fill=c, line=BORDER,
        bold=True, align=PP_ALIGN.CENTER)
    box(s, 1.30, y, 3.45, 0.78, t, size=9.5, fill=WHITE, line=BORDER,
        anchor=MSO_ANCHOR.MIDDLE)
    y += 0.78

secthead(s, 5.05, 2.78, 4.65, '期待される効果', size=10.5, h=0.32)
textbox(s, 5.05, 3.12, 4.65, 3.10,
        ['水害リスクを「じぶんごと」にできる：ハザードマップの数値では実感しにくい浸水の深さを、自宅前や通学路など見慣れた風景に実寸で重ねて示せる',
         '「自分の周辺は大丈夫」という思い込みを崩す：避難の必要性を認識しながら避難が遅れる事例への対策として、体感型の疑似体験が有効',
         '避難行動を事前に判断できる：「2階への垂直避難で足りるか、早めの立退き避難が必要か」を住まいごとに判断でき、マイ・タイムラインの検討材料になる',
         '防災教育・訓練に活用できる：屋内でも短時間で体験でき、学校・自主防災組織・町内会の訓練や出前講座で、子どもから高齢者まで直感的に理解できる',
         '平時のアプリ利用機会をつくる：災害時にしか開かれないという課題を補い、体験・共有を通じて認知とインストールが広がる。防災×としポの施策とも組み合わせやすい'],
        size=9.5, bullets=True)
textbox(s, L, 6.35, 9.40, 0.62,
        ['※ 実現方式イメージ：iOSはARKit、AndroidはARCore（Depth API）等を利用し、地面・床を検知して指定した深さの水面を描画。屋内・屋外いずれでも利用可能。',
         '※ 対応端末・表示精度、およびハザードマップの想定浸水深（洪水・高潮・津波）と連動した初期値表示の可否は、要件定義にて確定します。'],
        size=8.5, color=GRAY)

# ================================================================ P12 避難所チェックイン
s = S[11]
replace_paragraph(s, {'避難所チェックイン機能のイメージ': '機能拡張イメージ②｜避難所チェックイン（避難所受付システム）'})
for sh in s.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith('避難所の入り口に設置された'):
        sh.top = Inches(6.08); sh.width = Inches(4.30)
note(s, '※ 機能イメージであり、実装内容・マイナンバーカード連携の方式は要件定義にて確定します。', y=7.00, size=8.5)

# ================================================================ P13 各種防災機能
s = S[12]
old = find_shape(s, '機能紹介 ：')
if old is not None: old._element.getparent().remove(old._element)
bar = find_shape(s, 'ミニアプリは開発予定のためイメージです')
if bar is not None:
    bar.text_frame.paragraphs[0].runs[0].text = '※ 機能イメージであり、搭載する機能・遷移先は要件定義にて確定します。'
title(s, '機能拡張イメージ③｜ポータルアプリ上の防災機能')

# ================================================================ 重複スライド（P14）の削除
delete_slide(prs, 13)
S = list(prs.slides)

# ================================================================ P14 → 考えられるリスクと対策（防災アプリ向けに全面差し替え）
s = S[13]
clear(s)
title(s, '考えられるリスクと対策')
lead(s, '防災アプリの実装・移行にあたり想定されるリスクと、現時点で考えている対策。')
grid(s, L, 1.20, [1.85, 3.55, 4.00],
     ['リスク', '内容', '対策'],
     [['権利・引継ぎ', '「避難所へGo!」と同等機能を実装する際の著作権・利用許諾、データの引継ぎ',
       '現行ベンダーとの権利関係を事前に整理。仕様は要件として再定義し、必要に応じ許諾を取得'],
      ['災害時のアクセス集中', '発災直後に利用が集中し、表示や通知に遅延が生じる',
       '10万人の同時アクセス、10秒程度での一斉通知を前提に設計し、負荷試験を実施'],
      ['情報の正確性・鮮度', 'Lアラート・気象情報・避難所開設情報の反映漏れや遅延',
       '取得経路の冗長化と監視、取得失敗時のアラート。管理画面からの手動配信手段も確保'],
      ['通信断・権限未許可', '圏外や通信不安定、位置情報・通知の権限が未許可で機能が使えない',
       '地図をキャッシュしオフライン表示に対応。権限未許可でも地域選択により利用可能とする'],
      ['デジタル格差・多言語', '高齢者・外国人・障害のある方が利用しにくい',
       '多言語（日・英・中・韓）に対応し、文字サイズ等に配慮。防災ポータル等の代替手段も併用'],
      ['個人情報・セキュリティ', '位置情報や安否情報など、機微な情報を取り扱う',
       '広島市情報セキュリティポリシーに準拠。権限管理・暗号化・ログ保全。Pマーク／ISMS取得済'],
      ['関連市町との調整', '廿日市市・江田島市・熊野町でも現行アプリを利用中',
       '事前説明と移行方針の合意形成を先行して実施。周辺市町への拡大方針もあわせて協議'],
      ['移行・スケジュール', '現行アプリからの切替時期、並行運用の要否',
       '移行方式を協議のうえ並行運用期間を設定。リリース時期から逆算した工程管理を行う']],
     heights=[0.62]*8, size=9.5)

# ================================================================ P15 → 必要な情報・データと管理区分
s = S[14]
clear(s)
title(s, '必要な情報・データと管理区分')
lead(s, 'アプリの実装・運用に必要な情報とその提供元、システム／業務の管理区分について認識を合わせたい。')
grid(s, L, 1.20, [2.10, 3.45, 2.00, 1.85],
     ['情報・データ', '主な内容', '提供元', '管理区分'],
     [['避難所等情報', '名称・住所・標高・災害対応種別・収容人数・備蓄', '広島市', 'システム'],
      ['開設避難所情報', '避難所等の開設・閉鎖の状況', 'Lアラート', 'システム（自動取得）'],
      ['小学校区データ', '地域選択・エリア配信の区分（SHP形式）', '広島市', 'システム'],
      ['ハザード情報', '土砂・河川浸水・津波・高潮・洪水・雨水出水・ため池', '広島市／国・県', 'システム'],
      ['気象・防災情報', '気象警報・注意報、震源震度、各情報の解除報', 'Lアラート・気象庁', 'システム（自動取得）'],
      ['自治体からのお知らせ', '市からのお知らせ、イベント情報', '広島市', 'システム（管理画面）'],
      ['防災ハンドブック', '「たちまち防災」等のPDF', '広島市', 'システム（差替は管理画面）'],
      ['関連サイトのリンク', '広島市防災ポータル、web171・J-anpi 等', '広島市／各運営者', 'システム（URL設定）'],
      ['操作マニュアル・研修', '管理者・運用者向けのマニュアルと研修', 'FPM作成／広島市確認', '業務'],
      ['広報・利用促進', 'インストール案内、掲示物、キャラクターデザイン', '広島市・FPM', '業務']],
     heights=[0.44]*10, size=9.5)
note(s, '管理区分：システム＝アプリ・管理画面で管理／業務＝広島市およびFPMが業務として管理。'
        '提供元・更新頻度はヒアリングにて確認させていただきたい事項です。', y=6.10)

# ================================================================ P16 → 運用業務（バックオフィス）
s = S[15]
clear(s)
title(s, 'バックオフィスとして必要な業務')
lead(s, 'アプリの提供に加え、平時から災害時までを通じた運用業務が発生する。業務の認識と分担が合うか確認したい。')
grid(s, L, 1.20, [1.30, 4.35, 2.30, 1.45],
     ['フェーズ', '必要な業務', '実施主体', '頻度・時期'],
     [['事前設定', '避難所等情報・小学校区・ハザードレイヤの登録・更新', '危機管理対策課・FPM', '年次・随時'],
      ['事前設定', '通知カテゴリ・配信対象エリアの設定、非通知にできない情報の指定', '危機管理対策課', '年度初・随時'],
      ['事前設定', '職員アカウント・権限（管理者／運用者）の登録', '危機管理対策課・FPM', '随時'],
      ['平時運用', 'お知らせ・イベント情報の作成と配信', '危機管理対策課', '随時'],
      ['平時運用', '防災ハンドブック・リンク集の差替、利用促進の広報', '危機管理対策課', '随時'],
      ['警戒時', '避難情報の確認、Lアラート連携状況の監視', '危機管理対策課', '気象警報の発表時'],
      ['災害時', '緊急情報のプッシュ配信、エリアを指定した配信', '危機管理対策課', '発災時'],
      ['災害時', '開設避難所の反映確認、利用者からの問合せ対応', '危機管理対策課・FPM', '発災時'],
      ['事後', '配信実績・利用状況の確認と振り返り', '危機管理対策課', '災害後・月次'],
      ['保守', '障害対応、OS更新への対応、ストアの公開・更新', 'FPM', '随時・年次'],
      ['研修', '操作マニュアルに基づく職員研修の実施', 'FPM（実施）・危機管理対策課', '年次・体制変更時']],
     heights=[0.42]*11, size=9.5)
note(s, '※ 実施主体・頻度は現時点の想定です。運用フローの整備とあわせて、要求仕様の確定までに整理します。', y=6.36)

# ================================================================ 並び替え・ページ番号
reorder(prs, [0,1,2,3,4,5,6,7,8,10,11,12,13,14,15,9])
for i, sl in enumerate(prs.slides):
    if i == 0: continue
    pagenum(sl, i + 1)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print('saved', OUT, len(list(prs.slides)), 'slides')
